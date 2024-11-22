import re
from datetime import datetime
import tempfile
import mimetypes
import subprocess

import requests
from requests_futures.sessions import FuturesSession
from concurrent.futures import as_completed
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
from pptx import Presentation

from tipi_data.models.initiative import Initiative

from logger import get_logger
from extractors.config import LIMIT_DATE_TO_SYNC
from .initiatives_status import has_finished
from .api import ENDPOINT
from .legislative_period import LegislativePeriod
from .initiatives_attachments import (
    MIMETYPE_FILE_EXTENSIONS,
    ATTACHMENTS_WORKFLOW,
    get_current_phase,
    get_next_phase,
)


log = get_logger(__name__)


class InitiativesExtractor:
    def __init__(self):
        self.__legislative_period = LegislativePeriod().get()

    def __get_total(self):
        response = requests.get(ENDPOINT.format(method="proyecto/total"))
        if not response.ok:
            return 0
        return int(response.content)

    def extract(self):
        per_page = 20
        ROUNDING_UP = BOUNDARY = 1
        urls = [
            ENDPOINT.format(
                method="proyecto?offset={}&limit={}".format(offset, per_page)
            )
            for offset in range(
                1, (int(self.__get_total() / per_page)) + ROUNDING_UP + BOUNDARY
            )
        ]
        with FuturesSession() as session:
            futures = [session.get(url) for url in urls]
            for future in as_completed(futures):
                response = future.result()
                for initiative in response.json():
                    if not has_finished(initiative):
                        try:
                            self.__create_or_update(initiative)
                        except Exception as e:
                            log.warning(
                                "Expediente: {} - {}".format(
                                    str(initiative["idProyecto"]), e
                                )
                            )

    def __create_or_update(self, remote_initiative):
        try:
            initiative = Initiative.all.get(id=str(remote_initiative["idProyecto"]))
            if self.__too_old_to_process(initiative):
                return
        except Exception:
            initiative = Initiative()
        initiative["id"] = str(remote_initiative["idProyecto"])
        initiative["reference"] = str(remote_initiative["expedienteCamara"])
        initiative["title"] = remote_initiative["acapite"]
        initiative["initiative_type"] = remote_initiative["tipoProyecto"]
        initiative["history"] = [
            "{} ({})".format(
                remote_initiative["descripcionEtapa"],
                remote_initiative["descripcionSubEtapa"],
            )
        ]
        initiative["status"] = remote_initiative["estadoProyecto"]
        initiative["place"] = remote_initiative["origenProyecto"]
        initiative["url"] = remote_initiative["appURL"]
        initiative["created"] = self.__parse_date(
            remote_initiative["fechaIngresoExpediente"]
        )
        initiative["updated"] = self.__parse_date(
            remote_initiative["fechaIngresoExpediente"]
        )
        if "extra" not in initiative or initiative.extra == {}:
            initiative["extra"] = dict()
            initiative["extra"]["proponente"] = remote_initiative["iniciativa"]
            initiative["extra"]["ignored_attachments"] = list()
        self.__load_more_data(initiative)
        initiative.save()
        log.info("Iniciativa {} procesada".format(str(remote_initiative["idProyecto"])))

    def __load_more_data(self, initiative):
        response = requests.get(
            ENDPOINT.format(method="proyecto/{}/detalle".format(initiative["id"]))
        )
        self.__load_authors_from_response(initiative, response)
        self.__load_content_from_response(initiative, response)

    def __load_authors_from_response(self, initiative, response):
        if "listaAutores" in response.json().keys():
            initiative["author_deputies"] = [
                "{} {} [{}]".format(
                    author["nombres"].strip().title(),
                    author["apellidos"].strip().title(),
                    author["idParlamentario"],
                )
                for author in response.json()["listaAutores"]
            ]
            initiative["author_parliamentarygroups"] = list(
                {
                    author["partidoPolitico"]
                    for author in response.json()["listaAutores"]
                }
            )
        if "ministerios" in response.json().keys():
            initiative["author_others"] = response.json()["ministerios"]

    def __load_content_from_response(self, initiative, response):
        response = response.json()
        if "archivosAdjuntos" in response.keys():
            current_phase, current_phase_counter = get_current_phase(
                str(response["idProyecto"])
            )
            next_phase_index, next_phase = get_next_phase(current_phase)
            if next_phase_index != -1:
                attachments = response["archivosAdjuntos"]
                if current_phase_counter < len(
                    [a for a in attachments if a["infoAdjunto"] == current_phase]
                ):
                    next_phase_index = next_phase_index - 1
                for phase in ATTACHMENTS_WORKFLOW[next_phase_index:]:
                    self.__process_attachments_by_phase(
                        initiative,
                        [
                            attachment
                            for attachment in attachments
                            if attachment["infoAdjunto"] == phase
                        ],
                        phase,
                    )
        if not self.__has_content(initiative):
            self.__untag(initiative)
            initiative["content"] = [""]

    def __process_attachments_by_phase(self, initiative, attachments, phase):
        correct_counter = 0
        full_content = []
        for attachment in attachments:
            if attachment["idAdjunto"] in initiative["extra"]["ignored_attachments"]:
                continue
            response = requests.get(attachment["appURL"])
            if not response.ok:
                continue
            try:
                with tempfile.NamedTemporaryFile(
                    suffix=MIMETYPE_FILE_EXTENSIONS[attachment["tipoArchivo"]]
                ) as f:
                    f.write(bytes(response.content))
                    try:
                        mime_type, _ = mimetypes.guess_type(f.name)
                        if mime_type == "text/plain":
                            content = f.read().decode("utf-8").strip()
                        elif mime_type == "application/pdf":
                            content = extract_pdf_text(f.name).strip()
                        elif (
                            mime_type
                            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        ):
                            doc = Document(f)
                            content = "\n".join(
                                [para.text for para in doc.paragraphs]
                            ).strip()
                        elif mime_type == "application/msword":
                            result = subprocess.run(
                                ["antiword", f.name],
                                stdout=subprocess.PIPE,
                                stderr=subprocess.PIPE,
                            )
                            if result.returncode != 0:
                                raise Exception(
                                    f"Error al leer el archivo .doc: {result.stderr.decode('utf-8')}"
                                )
                            content = result.stdout.decode("utf-8").strip()
                        elif (
                            mime_type
                            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        ):
                            ppt = Presentation(f)
                            content = "\n".join(
                                [
                                    shape.text
                                    for slide in ppt.slides
                                    for shape in slide.shapes
                                    if hasattr(shape, "text")
                                ]
                            ).strip()
                        content = (
                            content.replace("\n", " ")
                            .replace("\f", " ")
                            .replace("\t", "")
                        )
                        content = [x for x in re.split(r"\.(?!\d)", content) if x != ""]
                        if len(content):
                            correct_counter = correct_counter + 1
                    except Exception:
                        content = [""]
                        initiative.extra["ignored_attachments"].append(
                            attachment["idAdjunto"]
                        )
                    f.close()
            except KeyError:
                pass  # Mimetype not found in our list
                content = [""]
                initiative.extra["ignored_attachments"].append(attachment["idAdjunto"])
            full_content = full_content + content
        if correct_counter:
            self.__untag(initiative)
            initiative["content"] = full_content
            initiative["extra"]["content_reference"] = phase
            initiative["extra"]["content_counter"] = correct_counter

    def __too_old_to_process(self, initiative):
        def frmt(dt):
            date_format = "%Y-%m-%d"
            return datetime.strptime(dt, date_format)

        return frmt(str(initiative["updated"]).split(" ")[0]) < frmt(LIMIT_DATE_TO_SYNC)

    def __has_content(self, initiative):
        try:
            saved_initiative = Initiative.all.get(id=initiative.id)
        except Exception:
            saved_initiative = dict()
        return "content" in saved_initiative or "content" in initiative

    def __untag(self, initiative):
        initiative["topics"] = []
        initiative["tags"] = []
        initiative["tagged"] = False

    def __parse_date(self, str_date):
        split_date = str_date.split("/")
        if len(split_date) != 3:
            return None
        return datetime(int(split_date[2]), int(split_date[1]), int(split_date[0]))
